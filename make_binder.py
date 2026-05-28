#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

C = {
    'dk':  RGBColor(0xC4, 0x5A, 0x1A),
    'md':  RGBColor(0xD9, 0x70, 0x36),
    'lt':  RGBColor(0xF0, 0x95, 0x6A),
    'pl':  RGBColor(0xFA, 0xE8, 0xDC),
    'cr':  RGBColor(0xFD, 0xF6, 0xF0),
    'bdk': RGBColor(0x2C, 0x1A, 0x0E),
    'bmd': RGBColor(0x5C, 0x33, 0x17),
    'wh':  RGBColor(0xFF, 0xFF, 0xFF),
    'gm':  RGBColor(0xCC, 0xCC, 0xCC),
    'gt':  RGBColor(0x66, 0x66, 0x66),
    'hdr': RGBColor(0xA0, 0x48, 0x12),
    'dc':  RGBColor(0x3C, 0x22, 0x10),
    'lx':  RGBColor(0xEE, 0xDD, 0xD0),
}
KR = 'Pretendard'
EN = 'Bernhard Modern'

prs = Presentation()
prs.slide_width  = Mm(210)
prs.slide_height = Mm(297)
BL = prs.slide_layouts[6]


def R(sl, x, y, w, h, bg=None, lc=None):
    s = sl.shapes.add_shape(1, Mm(x), Mm(y), Mm(w), Mm(h))
    if bg:
        s.fill.solid(); s.fill.fore_color.rgb = bg
    else:
        s.fill.background()
    if lc:
        s.line.color.rgb = lc; s.line.width = Pt(0.5)
    else:
        s.line.fill.background()
    return s


def T(sl, x, y, w, h, text, fn=KR, fs=10, fc=None, bold=False,
      align=PP_ALIGN.LEFT, bg=None, ml=2, mt=1.5, mr=2, mb=1.5):
    fc = fc or C['bdk']
    if bg:
        s = sl.shapes.add_shape(1, Mm(x), Mm(y), Mm(w), Mm(h))
        s.fill.solid(); s.fill.fore_color.rgb = bg
        s.line.fill.background()
    else:
        s = sl.shapes.add_textbox(Mm(x), Mm(y), Mm(w), Mm(h))
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = Mm(ml); tf.margin_top = Mm(mt)
    tf.margin_right = Mm(mr); tf.margin_bottom = Mm(mb)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = fn; r.font.size = Pt(fs)
    r.font.color.rgb = fc; r.font.bold = bold
    return s


def footer(sl, n):
    R(sl, 16, 283, 178, 0.4, bg=C['pl'])
    T(sl, 16, 284, 50, 7, '오블리브의원', KR, 7, C['gm'])
    T(sl, 76, 284, 68, 7, 'OBLIV CLINIC · DERMAWAVE', EN, 7, C['gm'], align=PP_ALIGN.CENTER)
    R(sl, 72, 286.5, 2, 2, bg=C['dk'])
    R(sl, 147, 286.5, 2, 2, bg=C['dk'])
    T(sl, 160, 284, 28, 7, f'0{n}', EN, 7, C['gm'], align=PP_ALIGN.RIGHT)


def hdr(sl, eyebrow, title, sub=''):
    R(sl, 0, 0, 210, 45, bg=C['dk'])
    ci = sl.shapes.add_shape(9, Mm(180), Mm(-10), Mm(26), Mm(26))
    ci.fill.solid(); ci.fill.fore_color.rgb = RGBColor(0xD0, 0x62, 0x20)
    ci.line.fill.background()
    T(sl, 16, 8, 178, 7, eyebrow, EN, 8, C['lt'])
    T(sl, 16, 17, 178, 13, title, KR, 18, C['wh'], bold=True)
    if sub:
        T(sl, 16, 32, 178, 10, sub, KR, 8, C['gm'])


def acc(sl):
    R(sl, 0, 0, 3.5, 297, bg=C['dk'])


def card(sl, y, title, dur, desc, tags):
    R(sl, 5, y, 200, 14, bg=C['dk'])
    T(sl, 8, y+3, 152, 9, title, KR, 10.5, C['wh'], bold=True)
    db = R(sl, 163, y+3, 40, 8, bg=C['hdr'])
    tf = db.text_frame; tf.margin_left=Mm(1); tf.margin_top=Mm(1)
    tf.margin_right=Mm(1); tf.margin_bottom=Mm(1)
    p = tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r = p.add_run(); r.text=dur; r.font.name=EN; r.font.size=Pt(8)
    r.font.color.rgb=C['wh']
    R(sl, 5, y+14, 200, 53, bg=C['cr'], lc=C['pl'])
    T(sl, 8, y+17, 194, 28, desc, KR, 8.5, C['gt'])
    tx = 8
    for tag in tags:
        tw = len(tag)*3.8+10
        ts = sl.shapes.add_shape(1, Mm(tx), Mm(y+47), Mm(tw), Mm(7))
        ts.fill.solid(); ts.fill.fore_color.rgb = C['wh']
        ts.line.color.rgb = C['md']; ts.line.width = Pt(0.5)
        ttf = ts.text_frame; ttf.margin_left=Mm(2); ttf.margin_top=Mm(1)
        ttf.margin_right=Mm(2); ttf.margin_bottom=Mm(1)
        tp = ttf.paragraphs[0]; tp.alignment=PP_ALIGN.CENTER
        tr = tp.add_run(); tr.text=tag
        tr.font.name=KR; tr.font.size=Pt(7.5); tr.font.color.rgb=C['dk']
        tx += tw+3


# ── SLIDE 1: COVER ──
s1 = prs.slides.add_slide(BL)
R(s1, 0, 0, 210, 297, bg=C['bdk'])
for ox, oy, ow, oc in [(140,-15,85,RGBColor(0x50,0x25,0x0C)),(-10,235,75,RGBColor(0x42,0x1C,0x07))]:
    ci = s1.shapes.add_shape(9, Mm(ox), Mm(oy), Mm(ow), Mm(ow))
    ci.fill.solid(); ci.fill.fore_color.rgb=oc; ci.line.fill.background()
bd = R(s1, 62, 79, 86, 9)
bd.line.color.rgb=C['md']; bd.line.width=Pt(0.6)
tf=bd.text_frame; tf.margin_left=Mm(2); tf.margin_top=Mm(1.5); tf.margin_right=Mm(2); tf.margin_bottom=Mm(1.5)
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
r=p.add_run(); r.text='CONSULTATION BINDER'; r.font.name=EN; r.font.size=Pt(7.5); r.font.color.rgb=C['md']
T(s1,20,93,170,22,'오블리브의원',KR,32,C['wh'],bold=True,align=PP_ALIGN.CENTER)
T(s1,20,116,170,9,'O B L I V   C L I N I C',EN,10,C['lt'],align=PP_ALIGN.CENTER)
R(s1,85,130,40,1,bg=C['md'])
T(s1,20,136,170,16,'더마웨이브 트리트먼트',KR,22,C['wh'],align=PP_ALIGN.CENTER)
T(s1,20,154,170,9,'DERMAWAVE TREATMENT GUIDE',EN,10,C['gm'],align=PP_ALIGN.CENTER)
for tag,tx in zip(['RF 고주파','근막 리모델링','비침습 시술','페이스 & 바디'],[17,61,105,149]):
    ts=s1.shapes.add_shape(1,Mm(tx),Mm(169),Mm(40),Mm(8))
    ts.fill.solid(); ts.fill.fore_color.rgb=C['dc']
    ts.line.color.rgb=RGBColor(0x7A,0x38,0x15); ts.line.width=Pt(0.5)
    ttf=ts.text_frame; ttf.margin_left=Mm(1); ttf.margin_top=Mm(1.5); ttf.margin_right=Mm(1); ttf.margin_bottom=Mm(1.5)
    tp=ttf.paragraphs[0]; tp.alignment=PP_ALIGN.CENTER
    tr=tp.add_run(); tr.text=tag; tr.font.name=KR; tr.font.size=Pt(8); tr.font.color.rgb=C['lt']
T(s1,20,277,170,7,'OBLIV CLINIC · DERMAWAVE · CONFIDENTIAL',EN,7,RGBColor(0x55,0x44,0x38),align=PP_ALIGN.CENTER)


# ── SLIDE 2: OVERVIEW ──
s2=prs.slides.add_slide(BL); acc(s2)
hdr(s2,'OVERVIEW','더마웨이브란?','고주파·저주파 복합 에너지와 그라스톤 괄사를 결합한 근막 리모델링 솔루션')
footer(s2,1)
R(s2,5,51,200,40,bg=C['dk'])
T(s2,8,54,194,9,'TUPS HL Frequency Dermawave',EN,12,C['wh'])
T(s2,8,64,194,25,
  'RF 고주파 심부열과 중저주파 전기 자극을 이용하여 피부와 전신을 관리하는 복합 솔루션입니다. '
  '괄사 도구를 이용해 근막을 리모델링하고 고주파·저주파 에너지로 세포 활성화를 촉진시켜 얼굴과 '
  '바디의 탄력을 근원적으로 개선합니다. 마취 없이 통증과 부작용이 매우 적어 누구나 편안하게 받을 수 있습니다.',
  KR,8.5,C['lx'])
T(s2,5,96,100,8,'핵심 특장점',KR,12,C['bmd'],bold=True)
R(s2,5,104,50,0.6,bg=C['dk'])
for (ft,fd),(cx,cy) in zip([
    ('피부 근원 케어','피부 세포 활성화 및 근막 리모델링으로 페이스·바디 문제의 근원부터 해결하여 몸의 밸런스를 되찾고 시술 효과를 장기간 유지합니다.'),
    ('스마트 복합 테라피','RET RF 주파 & 괄사와 LF 주파 & 글러브 전극을 동시에 사용하여 짧은 시간 안에 최대 관리 효율을 달성합니다.'),
    ('100% 비침습 시술','마취·통증 없이 진피층까지 에너지를 전달하여 탄력 개선, 리프팅, 셀룰라이트 분해, 통증 완화 효과를 제공합니다.'),
    ('신뢰성 높은 기기','국내 의료기기 제조사 생산, KC 안전인증 완료. 의료기기와 동일한 사양으로 기기 안정성 및 안전성을 확보하였습니다.'),
],[(5,108),(108,108),(5,183),(108,183)]):
    R(s2,cx,cy,96,68,bg=C['cr'],lc=C['pl'])
    R(s2,cx,cy,96,1.5,bg=C['dk'])
    T(s2,cx+3,cy+5,90,9,ft,KR,9.5,C['dk'],bold=True)
    R(s2,cx+3,cy+15,30,0.4,bg=C['md'])
    T(s2,cx+3,cy+19,90,46,fd,KR,8,C['gt'])


# ── SLIDE 3: MECHANISM ──
s3=prs.slides.add_slide(BL); acc(s3)
hdr(s3,'MECHANISM','시술 원리','고주파(RF) + 저주파(LF) + 그라스톤 괄사의 3중 복합 작용')
footer(s3,2)
for (mt,md),mx in zip([
    ('RF 고주파 (모노폴라)','2~9 MHz 심부열 에너지로 진피층·SMAS층에 도달하여 콜라겐·엘라스틴 재생을 촉진합니다.'),
    ('저주파 (LF) 전기 자극','중저주파 전기 자극으로 근육을 선택적으로 수축시키고 신경·근막 기능을 활성화합니다.'),
    ('그라스톤 괄사 도구','스테인리스 스틸 기구로 유착된 근막을 물리적으로 이완하여 조직 기능을 근원적으로 개선합니다.'),
],[5,72,139]):
    R(s3,mx,51,63,55,bg=C['cr'],lc=C['pl'])
    R(s3,mx,51,63,1.5,bg=C['dk'])
    T(s3,mx+3,56,57,9,mt,KR,9,C['dk'],bold=True)
    T(s3,mx+3,67,57,36,md,KR,8,C['gt'])
R(s3,5,112,200,0.4,bg=C['pl'])
T(s3,5,117,100,8,'핸드피스 타입 안내',KR,12,C['bmd'],bold=True)
R(s3,5,125,65,0.6,bg=C['dk'])
hy=130
for hl,hst,hd in [
    ('D2 TYPE','물고기형 — 페이스·데콜테·손·발·다리','얼굴 윤곽 정리, 부종 완화, 리프팅에 특화된 타입으로 섬세한 부위에 정확하게 적용합니다.'),
    ('D3 TYPE','반달형 — 승모근·어깨·팔·복부·다리','슬리밍 및 림프 순환, 근막 이완에 최적화된 타입으로 중간 크기 부위를 효과적으로 커버합니다.'),
    ('D4 TYPE','뱀부형 — 승모근·등·옆구리·어깨·다리','넓은 근막 범위의 림프 순환과 근막 이완에 특화된 타입으로 광범위 부위를 효율적으로 관리합니다.'),
]:
    lb=R(s3,5,hy,18,8,bg=C['dk'])
    tf=lb.text_frame; tf.margin_left=Mm(1); tf.margin_top=Mm(1); tf.margin_right=Mm(1); tf.margin_bottom=Mm(1)
    lp=tf.paragraphs[0]; lp.alignment=PP_ALIGN.CENTER
    lr=lp.add_run(); lr.text=hl; lr.font.name=EN; lr.font.size=Pt(7); lr.font.color.rgb=C['wh']; lr.font.bold=True
    T(s3,26,hy,178,8,hst,KR,9,C['bmd'],bold=True)
    T(s3,26,hy+9,178,15,hd,KR,8,C['gt'])
    R(s3,5,hy+27,200,0.3,bg=C['pl'])
    hy+=30
R(s3,5,hy+3,200,22,bg=C['dk'])
T(s3,8,hy+6,194,8,'기대 효과',KR,11,C['wh'],bold=True)
T(s3,8,hy+15,194,9,'건강한 근막 형성  ·  슬리밍 / 쉐이핑  ·  주름 개선  ·  체중 감량 증진  ·  부종 / 셀룰라이트 완화',KR,8.5,C['lx'],align=PP_ALIGN.CENTER)


# ── SLIDE 4: TARGET ──
s4=prs.slides.add_slide(BL); acc(s4)
hdr(s4,'RECOMMENDATION','이런 분께 추천합니다','더마웨이브 트리트먼트 추천 고객 유형')
footer(s4,3)
T(s4,5,50,80,8,'추천 대상',KR,12,C['bmd'],bold=True)
R(s4,5,58,45,0.6,bg=C['dk'])
ty=62
for tg in ['수술 없이 자연스러운 탄력과 리프팅을 원하시는 분','1회 관리만으로도 즉각적인 효과가 필요하신 분',
           '일상생활에 지장 없는 관리를 원하시는 분','신체 순환 저하로 붓기가 잘 빠지지 않으시는 분',
           '몸에 에너지가 부족하고 만성 피로에 시달리는 분','셀룰라이트 완화 및 군살 정리가 필요하신 분',
           '피부과 시술 이후 빠른 회복이 필요하신 분','다이어트로 인한 탄력 저하가 고민이신 분']:
    R(s4,5,ty+3,2.5,2.5,bg=C['dk'])
    T(s4,10,ty,88,13,tg,KR,8.5,C['bdk'])
    R(s4,5,ty+14,93,0.3,bg=C['pl'])
    ty+=16
T(s4,107,50,95,8,'관리 가능 부위',KR,12,C['bmd'],bold=True)
R(s4,107,58,58,0.6,bg=C['dk'])
ay=62
for at,ad in [('페이스 & 데콜테','얼굴 리프팅 · 쉐이핑 · 안색 개선 · 주름 완화'),
              ('상체 전면','복부 슬리밍 · 골반 교정 · 부유방 제거 · 팔뚝라인 정리'),
              ('상체 후면','척추(거북목) 교정 · 승모근 완화 · 등 군살 정리'),
              ('하체','다리 부종 완화 · 슬리밍 · 셀룰라이트 제거 · 피부 토닝')]:
    R(s4,107,ay,98,48,bg=C['cr'],lc=C['pl'])
    R(s4,107,ay,98,1.5,bg=C['dk'])
    T(s4,110,ay+6,92,9,at,KR,9.5,C['dk'],bold=True)
    T(s4,110,ay+17,92,26,ad,KR,8,C['gt'])
    ay+=52


# ── SLIDE 5: TREATMENTS 1 ──
s5=prs.slides.add_slide(BL); acc(s5)
hdr(s5,'TREATMENT MENU','트리트먼트 프로그램','부위별 맞춤 시술 — 관리 효과 및 시간 안내')
footer(s5,4)
ty=50
for tt,td,tdesc,ttags in [
    ('페이스 & 데콜테 트리트먼트','약 20분',
     '고주파를 통해 진피층의 콜라겐과 탄력섬유를 활성화시키고 림프에 쌓인 노폐물을 배출시켜 얼굴 라인을 매끄럽게 정리합니다. 입체적인 볼륨감을 부여하여 자연스러운 안티에이징 효과를 완성합니다.',
     ['페이스 리프팅','쉐이핑','안색 개선','볼륨 부여']),
    ('상체 전면 트리트먼트','약 30분',
     '체내에 쌓인 냉기를 배출하고 코어의 힘을 길러 장기적으로 면역력을 올리며 셀룰라이트를 감소시켜 건강한 신체를 가꿉니다. 디톡스·복부 슬리밍·부유방 제거·팔뚝라인 정리에 효과적입니다.',
     ['복부 슬리밍','골반 교정','부유방 제거','팔뚝라인 정리','면역력 증강']),
    ('상체 후면 트리트먼트','약 30분',
     '기립근을 강화하고 틀어진 몸의 중심을 바로잡습니다. 만성적인 어깨결림과 승모근 솟음을 완화시켜 신체 전반에 유연함과 편안함을 주고 라인을 정리하여 한결 가벼운 신체를 가꿉니다.',
     ['척추 교정','승모근 완화','라운드숄더 개선','군살 정리']),
]:
    card(s5,ty,tt,td,tdesc,ttags); ty+=74


# ── SLIDE 6: TREATMENTS 2 ──
s6=prs.slides.add_slide(BL); acc(s6)
hdr(s6,'TREATMENT MENU','트리트먼트 프로그램','하체 · 전신 — 관리 효과 및 시간 안내')
footer(s6,5)
ty=50
for tt,td,tdesc,ttags in [
    ('하체 트리트먼트','약 40분',
     '고착된 근막을 풀어주고 림프의 순환을 도와 붓기를 제거합니다. 비대해진 셀룰라이트를 감소시켜 보다 가볍고 매끈해진 하체를 가꾸는 데 도움을 줍니다.',
     ['다리 부종 완화','슬리밍','셀룰라이트 제거','피부 토닝','착색 완화']),
    ('전신 트리트먼트','약 90분',
     '근막·림프·혈관·셀룰라이트를 모두 자극하여 유착된 근막을 재배열하고 피하지방층의 지방세포를 감소시킵니다. 림프와 혈관의 흐름을 개선시켜 체온을 높이고 염증을 완화하여 몸속부터 근원적인 변화를 일으킵니다.',
     ['신체 순환 개선','전신 슬리밍','신체 에너지 강화','근막 재배열','염증 완화']),
]:
    card(s6,ty,tt,td,tdesc,ttags); ty+=74
R(s6,5,ty+3,200,0.4,bg=C['pl'])
R(s6,5,ty+8,200,38,bg=C['dk'])
T(s6,8,ty+11,194,9,'트리트먼트 요약',KR,11,C['wh'],bold=True)
T(s6,8,ty+22,194,11,'페이스&데콜테 20분  ·  상체전면 30분  ·  상체후면 30분  ·  하체 40분  ·  전신 90분',KR,8.5,C['lx'])
T(s6,8,ty+31,194,11,'효과의 지속을 위해 주 2~3회 시술, 10회 이상을 권장드립니다.',KR,8.5,C['lx'])


# ── SLIDE 7: Q&A ──
s7=prs.slides.add_slide(BL); acc(s7)
hdr(s7,'Q & A','자주 묻는 질문','상담 시 환자분들이 가장 많이 문의하시는 내용을 정리했습니다')
footer(s7,6)
qy=50
for qq,qa in [
    ('Q. 시술 시간은 얼마나 걸리나요?',
     '일반적으로 부위당 15~30분 소요됩니다. 하나의 장비로 고주파·MFR 툴·저주파·매직글러브 시술이 모두 가능하므로 기존 대비 짧은 시간 안에 복합 시술을 받으실 수 있습니다.'),
    ('Q. 시술 효과는 바로 나타나나요?',
     '네, 1회 시술만으로도 즉각적인 효과를 확인하실 수 있습니다. 효과의 지속과 극대화를 위해 주 2~3회, 총 10회 이상의 시술을 권장드립니다.'),
    ('Q. 통증이 있거나 일상생활에 지장이 있나요?',
     '일반 고주파와 달리 더마웨이브는 심부열 에너지를 활용하기 때문에 통증 없이 부드럽고 편안한 관리가 가능합니다. 시술 직후 바로 일상생활로 복귀하실 수 있습니다.'),
    ('Q. 부작용은 없나요?',
     '더마웨이브는 비침습적 시술이며 신체의 자연치유능력을 극대화하는 방식으로 부작용은 거의 없습니다. 다만 근막 유착이 오래된 부위는 유착 해소 과정에서 피부 표면이 약간 붉어질 수 있으나 자연스럽게 사라집니다.'),
]:
    R(s7,5,qy,2.5,38,bg=C['dk'])
    T(s7,10,qy+3,194,9,qq,KR,9.5,C['dk'],bold=True)
    T(s7,10,qy+14,194,22,qa,KR,8.5,C['gt'])
    R(s7,5,qy+40,200,0.3,bg=C['pl'])
    qy+=45
R(s7,5,qy+4,200,35,bg=C['dk'])
T(s7,8,qy+7,194,9,'상담 안내',KR,11,C['wh'],bold=True)
T(s7,8,qy+18,194,20,
  '담당 의료진과의 1:1 상담을 통해 고객님의 체형·피부 상태에 최적화된 트리트먼트 플랜을 제안해 드립니다. '
  '궁금하신 사항은 언제든지 오블리브의원 상담팀으로 문의해 주시기 바랍니다.',
  KR,8.5,C['lx'])

prs.save('/home/user/Marketing/obliv_dermawave_binder.pptx')
print('Done.')
